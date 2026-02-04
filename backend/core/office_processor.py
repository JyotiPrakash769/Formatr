import os
from pptx import Presentation
from PIL import Image
import io

class OfficeProcessor:
    @staticmethod
    def extract_pptx_images(input_path: str, output_dir: str) -> str:
        """Extract all images from PowerPoint file"""
        prs = Presentation(input_path)
        filename = os.path.basename(input_path)
        name, _ = os.path.splitext(filename)
        
        # Create subfolder for images
        images_folder = os.path.join(output_dir, f"{name}_images")
        os.makedirs(images_folder, exist_ok=True)
        
        image_count = 0
        
        # Iterate through slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Check all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    # Extract image
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Determine extension
                    ext = image.ext
                    image_filename = f"slide{slide_num}_img{image_count + 1}.{ext}"
                    image_path = os.path.join(images_folder, image_filename)
                    
                    # Save image
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    
                    image_count += 1
        
        # Create a summary file
        summary_path = os.path.join(images_folder, "summary.txt")
        with open(summary_path, 'w') as f:
            f.write(f"Extracted {image_count} images from {filename}\n")
            f.write(f"Total slides: {len(prs.slides)}\n")
        
        return images_folder
    
    @staticmethod
    def extract_pptx_text(input_path: str, output_dir: str) -> str:
        """Extract all text from PowerPoint file"""
        prs = Presentation(input_path)
        filename = os.path.basename(input_path)
        name, _ = os.path.splitext(filename)
        output_path = os.path.join(output_dir, f"{name}_text.txt")
        
        with open(output_path, 'w', encoding='utf-8') as f:
            for slide_num, slide in enumerate(prs.slides, 1):
                f.write(f"=== SLIDE {slide_num} ===\n\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            f.write(f"{text}\n\n")
                
                f.write("\n" + "="*50 + "\n\n")
        
        return output_path
    
    @staticmethod
    def get_pptx_info(input_path: str, output_dir: str) -> str:
        """Get metadata and information about PowerPoint file"""
        prs = Presentation(input_path)
        filename = os.path.basename(input_path)
        name, _ = os.path.splitext(filename)
        output_path = os.path.join(output_dir, f"{name}_info.txt")
        
        # Count images
        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_count += 1
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"PowerPoint File Information\n")
            f.write(f"{'='*50}\n\n")
            f.write(f"Filename: {filename}\n")
            f.write(f"Total Slides: {len(prs.slides)}\n")
            f.write(f"Total Images: {image_count}\n")
            f.write(f"Slide Size: {prs.slide_width} x {prs.slide_height}\n")
            
            # Core properties if available
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                f.write(f"Title: {prs.core_properties.title}\n")
            if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
                f.write(f"Author: {prs.core_properties.author}\n")
            if hasattr(prs.core_properties, 'subject') and prs.core_properties.subject:
                f.write(f"Subject: {prs.core_properties.subject}\n")
        
        return output_path
